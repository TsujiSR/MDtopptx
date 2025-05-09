"""
md2pptx-builder - PowerPoint builder
"""

import os
import logging
from typing import List, Dict, Any, Optional, Tuple, Union
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from md2pptx_builder.utils import is_valid_image, get_image_dimensions

logger = logging.getLogger(__name__)

class PPTXBuilder:
    """MarkdownからPowerPointを生成するクラス"""
    
    def __init__(self, 
                 background_path: str, 
                 logo_path: str, 
                 template_path: Optional[str] = None,
                 font_family: str = "メイリオ",
                 verbose: bool = False):
        """
        Args:
            background_path: 背景画像のパス
            logo_path: ロゴ画像のパス
            template_path: テンプレートPPTXのパス（オプション）
            font_family: 使用するフォント
            verbose: 詳細ログを出力するかどうか
        """
        self.background_path = background_path
        self.logo_path = logo_path
        self.template_path = template_path
        self.font_family = font_family
        self.verbose = verbose
        
        # フォント設定の英語フォールバック対応
        self.fallback_font = "Arial"
        if "明朝" in self.font_family or "Serif" in self.font_family:
            self.fallback_font = "Times New Roman"
        
        # 画像ファイルのチェック
        if not is_valid_image(background_path):
            raise ValueError(f"無効な背景画像: {background_path}")
        if not is_valid_image(logo_path):
            raise ValueError(f"無効なロゴ画像: {logo_path}")
        
        # プレゼンテーション作成
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
            logger.info(f"テンプレートを使用: {template_path}")
        else:
            self.prs = Presentation()
            logger.info("新規プレゼンテーションを作成")
            
        # デフォルトのスライドサイズを16:9に設定（テンプレートが無い場合）
        if not template_path:
            self.prs.slide_width = Inches(16 * 0.75)  # 16:9 比率
            self.prs.slide_height = Inches(9 * 0.75)
    
    def create_slide(self, slide_data: Dict[str, Any], total_slides: int) -> None:
        """スライドを作成する
        
        Args:
            slide_data: スライドデータ（タイトル、コンテンツなど）
            total_slides: スライドの総数
        """
        # レイアウトインデックス6は白紙のスライド
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        # 背景画像設定
        self._apply_background(slide)
        
        # ロゴ設定
        self._add_logo(slide)
        
        # タイトル追加
        title = slide_data.get("title", f"スライド {slide_data['index'] + 1}")
        self._add_title(slide, title)
        
        # コンテンツ追加
        self._add_content(slide, slide_data["content"])
        
        # スライド番号追加
        current_slide = slide_data["index"] + 1
        self._add_slide_number(slide, current_slide, total_slides)
        
        logger.info(f"スライド {current_slide}/{total_slides} を作成: {title}")
    
    def _apply_background(self, slide) -> None:
        """スライドに背景画像を適用する
        
        Args:
            slide: スライドオブジェクト
        """
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 白背景
        
        try:
            # 背景画像を全面に設定
            slide.shapes.add_picture(
                self.background_path,
                0, 0,
                width=self.prs.slide_width,
                height=self.prs.slide_height
            )
        except Exception as e:
            logger.error(f"背景画像の適用に失敗: {e}")
    
    def _add_logo(self, slide) -> None:
        """スライドにロゴを追加する
        
        Args:
            slide: スライドオブジェクト
        """
        try:
            # ロゴを右上に配置
            logo_width = Inches(1.2)  # ロゴサイズを調整 (1.5→1.2)
            logo = slide.shapes.add_picture(
                self.logo_path,
                self.prs.slide_width - logo_width - Inches(0.4),  # 右マージン (0.25→0.4)
                Inches(0.4),  # 上マージン (0.25→0.4)
                width=logo_width
            )
            if self.verbose:
                logger.debug(f"ロゴを追加: {logo.width} x {logo.height}")
        except Exception as e:
            logger.error(f"ロゴの追加に失敗: {e}")
    
    def _add_title(self, slide, title: str) -> None:
        """スライドにタイトルを追加する
        
        Args:
            slide: スライドオブジェクト
            title: タイトルテキスト
        """
        title_box = slide.shapes.add_textbox(
            Inches(0.5),  # 左マージン縮小 (0.6→0.5)
            Inches(0.5),  # 上マージン縮小 (0.6→0.5)
            self.prs.slide_width - Inches(2.0),  # 幅
            Inches(1.0)  # 高さ縮小 (1.2→1.0)
        )
        
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # 左揃え
        
        # タイトル下の余白を追加
        title_frame.paragraphs[0].space_after = Pt(4)  # 余白縮小 (6→4)
        
        title_run = title_frame.paragraphs[0].runs[0]
        title_run.font.size = Pt(30)  # タイトルサイズ変更 (36→30)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(0, 0, 0)  # 黒色
        title_run.font.name = self.font_family
        title_run.font.name_ascii = self.fallback_font  # 英文用フォールバック
    
    def _add_content(self, slide, content_ast: List[Dict[str, Any]]) -> None:
        """スライドにMarkdownコンテンツを追加する
        
        Args:
            slide: スライドオブジェクト
            content_ast: コンテンツのAST
        """
        # デバッグログ - AST構造を詳細に出力
        import json
        if logger.isEnabledFor(logging.INFO):
            logger.info(f"コンテンツAST: {json.dumps(content_ast, indent=2, ensure_ascii=False)}")
        
        # コンテンツ領域の定義 - マージン改善
        content_box = slide.shapes.add_textbox(
            Inches(0.7),  # 左マージン (1.0→0.7)
            Inches(1.6),  # タイトル下の開始位置 (2.0→1.6)
            self.prs.slide_width - Inches(1.5),  # 幅（右マージン縮小 2.0→1.5）
            self.prs.slide_height - Inches(2.0)  # 高さ（拡大 2.5→2.0）
        )
        
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = True  # テキストに合わせて自動調整
        
        # 段落間のスペーシング設定
        text_frame.paragraphs[0].space_after = Pt(4)  # 段落後の間隔縮小（8→4）
        
        # 最初の段落をクリア
        if text_frame.paragraphs:
            first_paragraph = text_frame.paragraphs[0]
            first_paragraph.text = ""
            current_paragraph = first_paragraph
        else:
            current_paragraph = text_frame.add_paragraph()
        
        # より直接的なアプローチでASTを処理
        for node in content_ast:
            self._process_node_direct(node, text_frame)
    
    def _process_node_direct(self, node: Dict[str, Any], text_frame) -> None:
        """ASTノードを直接処理して段落に変換する
        
        Args:
            node: ASTノード
            text_frame: テキストフレーム
        """
        node_type = node.get("type", "")
        
        if node_type == "blank_line":
            # 空行は段落を追加して空白を作る
            p = text_frame.add_paragraph()
            p.space_after = Pt(4)  # 空行の後の間隔（8→4）
            
        elif node_type == "paragraph":
            # 段落テキストの抽出と追加
            p = text_frame.add_paragraph()
            p.space_after = Pt(8)  # 段落間の間隔（12→8）
            
            children = node.get("children", [])
            for child in children:
                self._process_inline_node(child, p)
            
        elif node_type == "heading":
            # 見出しの処理
            level = node.get("attrs", {}).get("level", 2)
            children = node.get("children", [])
            
            p = text_frame.add_paragraph()
            p.space_before = Pt(10)  # 見出し前の間隔（16→10）
            p.space_after = Pt(6)  # 見出し後の間隔（8→6）
            
            # 見出しレベルに応じたフォントサイズ設定
            font_size = {
                2: 26,  # H2（28→26）
                3: 22,  # H3（24→22）
                4: 18,  # H4（20→18）
                5: 16,  # H5（18→16）
                6: 15   # H6（16→15）
            }.get(level, 26)
            
            for child in children:
                run = p.add_run()
                run.text = self._node_to_text(child)
                run.font.bold = True
                run.font.size = Pt(font_size)
                run.font.name = self.font_family
                run.font.name_ascii = self.fallback_font
            
        elif node_type == "list":
            # リストの処理
            # リスト前に少し余白を追加（小さめに調整）
            space_before = text_frame.add_paragraph()
            space_before.space_after = Pt(2)  # 余白縮小（4→2）
            
            # リスト処理
            self._process_list_direct(node, text_frame)
            
            # リスト後に余白を追加（小さめに調整）
            space_after = text_frame.add_paragraph()
            space_after.space_after = Pt(4)  # 余白縮小（8→4）
        
        elif node_type == "block_code":
            # コードブロックの処理
            code_text = node.get("raw", "")
            lang = node.get("attrs", {}).get("info", "")
            
            p = text_frame.add_paragraph()
            p.space_before = Pt(12)  # コードブロック前の間隔（6→12）
            p.space_after = Pt(12)  # コードブロック後の間隔（6→12）
            
            # 言語情報があれば追加
            if lang:
                lang_run = p.add_run()
                lang_run.text = f"{lang}:\n"
                lang_run.font.bold = True
                lang_run.font.size = Pt(14)
                self._apply_font_to_run(lang_run)
            
            # コードブロックの追加
            code_run = p.add_run()
            code_run.text = code_text
            code_run.font.size = Pt(14)
            code_run.font.name = "Consolas"  # コード用モノスペースフォント
            code_run.font.name_ascii = "Consolas"
    
    def _process_list_direct(self, node: Dict[str, Any], text_frame) -> None:
        """リストを処理してテキストフレームに追加する
        
        Args:
            node: リストノード
            text_frame: 追加先のテキストフレーム
        """
        is_ordered = node.get("attrs", {}).get("ordered", False)
        list_items = node.get("children", [])
        depth = node.get("attrs", {}).get("depth", 0)
        
        for i, item in enumerate(list_items):
            try:
                # リスト項目テキストを抽出（改善版）
                item_text = self._extract_list_item_text(item)
                
                logger.info(f"リスト項目 {i} テキスト: '{item_text}'")
                
                # 空のリスト項目は処理しない
                if not item_text.strip():
                    logger.info(f"空のリスト項目をスキップ: {item}")
                    continue
                
                # コロンの後にスペースを追加（必要な場合）
                if ":" in item_text and " " not in item_text.split(":", 1)[0]:
                    parts = item_text.split(":", 1)
                    if len(parts) == 2:
                        item_text = f"{parts[0]}: {parts[1].lstrip()}"
                        logger.info(f"コロン処理後: '{item_text}'")
                
                # リスト番号またはマーカーを生成
                marker = f"{i+1}." if is_ordered else "•"
                
                # リスト項目を段落として追加
                p = text_frame.add_paragraph()
                p.level = depth  # インデントレベル
                p.space_after = Pt(3)  # 項目間の間隔を減らす（5→3）
                
                run = p.add_run()
                # マーカーとテキストを結合
                run.text = f"{marker} {item_text}"
                
                # フォントスタイル設定 - サイズをやや大きく
                base_size = 18 if depth == 0 else 16  # 基本サイズをやや小さく調整
                run.font.size = Pt(base_size - depth)  # ネストレベルに応じて小さく
                self._apply_font_to_run(run)
                
                # 子リストがあれば再帰的に処理
                self._process_nested_lists(item, text_frame)
                
            except Exception as e:
                logger.error(f"リスト項目処理エラー: {e}", exc_info=True)
                # エラーがあっても処理を続行
    
    def _process_nested_lists(self, parent_item: Dict[str, Any], text_frame) -> None:
        """ネストされたリストを処理する
        
        Args:
            parent_item: 親リスト項目
            text_frame: テキストフレーム
        """
        children = parent_item.get("children", [])
        
        # children内のリストを検索
        for child in children:
            if child.get("type") == "list":
                # 子リストがある場合、深さを増やして処理
                child_list = child
                depth = child_list.get("attrs", {}).get("depth", 0) + 1
                
                # 深さ情報を設定
                if "attrs" not in child_list:
                    child_list["attrs"] = {}
                child_list["attrs"]["depth"] = depth
                
                # リストを処理
                self._process_list_direct(child_list, text_frame)
                
                # ネストされたリスト後に余分な空白を追加（見やすさのため）
                space_para = text_frame.add_paragraph()
                space_para.space_after = Pt(4)
                
                # ネストされたリストのマーカーを追加
                if child_list.get("attrs", {}).get("ordered", False):
                    marker = f"{depth}. "
                else:
                    marker = "• "
                
                # マーカーを段落に追加
                marker_run = space_para.add_run()
                marker_run.text = marker
                marker_run.font.size = Pt(14)
                self._apply_font_to_run(marker_run)
    
    def _node_to_text(self, node: Dict[str, Any]) -> str:
        """ノードからテキストを抽出する
        
        Args:
            node: ASTノード
            
        Returns:
            str: 抽出されたテキスト
        """
        if "children" not in node:
            return node.get("raw", node.get("text", ""))
        
        texts = []
        for child in node["children"]:
            child_type = child.get("type", "")
            
            if child_type == "text":
                texts.append(child.get("raw", child.get("text", "")))
            elif child_type == "strong":
                texts.append("**" + self._node_to_text(child) + "**")
            elif child_type == "emphasis":
                texts.append("*" + self._node_to_text(child) + "*")
            elif child_type == "code":
                texts.append("`" + child.get("raw", child.get("text", "")) + "`")
            elif "children" in child:
                texts.append(self._node_to_text(child))
        
        return "".join(texts)
    
    def _list_item_to_text(self, item: Dict[str, Any]) -> str:
        """リスト項目からテキストを抽出する
        
        Args:
            item: リスト項目ノード
            
        Returns:
            str: 抽出されたテキスト
        """
        texts = []
        if logger.isEnabledFor(logging.INFO):
            logger.info(f"処理中のリスト項目: {item}")
        
        for child in item.get("children", []):
            if logger.isEnabledFor(logging.INFO):
                logger.info(f"子要素タイプ: {child.get('type')}")
            
            if child.get("type") == "list":
                # ネストされたリストはスキップ
                continue
                
            elif child.get("type") == "block_text":
                # block_text要素を処理（mistune 3.1.3の特殊構造）
                if logger.isEnabledFor(logging.INFO):
                    logger.info(f"block_text要素処理: {child}")
                
                # 直接block_textの子要素または文字列を抽出
                if "children" in child:
                    for block_child in child.get("children", []):
                        if block_child.get("type") == "text":
                            text = block_child.get("raw", block_child.get("text", ""))
                            texts.append(text)
                            if logger.isEnabledFor(logging.INFO):
                                logger.info(f"抽出されたテキスト (block_text内): {text}")
                        elif block_child.get("type") == "strong":
                            strong_text = "**" + self._node_to_text(block_child) + "**"
                            texts.append(strong_text)
                        elif block_child.get("type") == "emphasis":
                            emph_text = "*" + self._node_to_text(block_child) + "*"
                            texts.append(emph_text)
                        # 他のインラインタイプも考慮
                        elif "children" in block_child:
                            texts.append(self._node_to_text(block_child))
                else:
                    # 子要素がない場合は直接テキストを抽出
                    text = child.get("raw", child.get("text", ""))
                    if text:
                        texts.append(text)
                
            elif child.get("type") == "paragraph":
                # 段落から各要素を抽出
                para_text = []
                
                for para_child in child.get("children", []):
                    child_type = para_child.get("type", "")
                    
                    if child_type == "text":
                        para_text.append(para_child.get("raw", para_child.get("text", "")))
                    elif child_type == "strong":
                        # 強調表示はマーカーつきで抽出（後で処理）
                        strong_children = para_child.get("children", [])
                        strong_text = "".join(c.get("raw", c.get("text", "")) for c in strong_children)
                        para_text.append(f"**{strong_text}**")
                    elif child_type == "emphasis":
                        # 斜体もマーカーつきで抽出
                        emph_children = para_child.get("children", [])
                        emph_text = "".join(c.get("raw", c.get("text", "")) for c in emph_children)
                        para_text.append(f"*{emph_text}*")
                    elif "children" in para_child:
                        # その他の複合要素
                        para_text.append(self._node_to_text(para_child))
                
                texts.append("".join(para_text))
                
            elif child.get("type") == "text":
                texts.append(child.get("raw", child.get("text", "")))
            
            elif child.get("type") == "strong":
                strong_text = self._node_to_text(child)
                texts.append(f"**{strong_text}**")
                
            elif child.get("type") == "emphasis":
                emph_text = self._node_to_text(child)
                texts.append(f"*{emph_text}*")
            
            # その他のノードタイプも対応
            elif "children" in child:
                texts.append(self._node_to_text(child))
        
        result = "".join(texts)
        if logger.isEnabledFor(logging.INFO):
            logger.info(f"リスト項目から抽出された最終テキスト: {result}")
        return result
    
    def _add_slide_number(self, slide, current: int, total: int) -> None:
        """スライド番号を追加する
        
        Args:
            slide: スライドオブジェクト
            current: 現在のスライド番号
            total: スライドの総数
        """
        # フッター領域にスライド番号を配置
        number_box = slide.shapes.add_textbox(
            self.prs.slide_width - Inches(1.5),  # 右マージン（2→1.5に調整）
            self.prs.slide_height - Inches(0.6),  # 下部に配置（0.8→0.6に調整）
            Inches(1.0),  # 幅
            Inches(0.3)   # 高さ
        )
        
        number_frame = number_box.text_frame
        number_frame.text = f"{current}/{total}"
        number_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        number_run = number_frame.paragraphs[0].runs[0]
        number_run.font.size = Pt(12)  # フォントサイズ小さく（14→12pt）
        number_run.font.color.rgb = RGBColor(80, 80, 80)  # グレー
        self._apply_font_to_run(number_run)
    
    def _apply_font_to_run(self, run) -> None:
        """テキストランにフォント設定を適用する
        
        Args:
            run: テキストラン
        """
        try:
            # 日本語フォント設定
            run.font.name = self.font_family
            # 英語フォールバック設定
            run.font.name_ascii = self.fallback_font
            # Eastasia用フォント（CJK文字）も設定
            if hasattr(run.font, 'name_eastasia'):
                run.font.name_eastasia = self.font_family
        except Exception as e:
            logger.warning(f"フォント適用エラー: {e}")
    
    def _process_inline_node(self, node: Dict[str, Any], paragraph) -> None:
        """インラインノードを処理して段落に追加する
        
        Args:
            node: インラインノードデータ
            paragraph: 追加先の段落
        """
        node_type = node.get("type", "")
        
        if node_type == "text":
            run = paragraph.add_run()
            run.text = node.get("raw", "")
            run.font.size = Pt(18)  # 本文フォントサイズ
            self._apply_font_to_run(run)
            
        elif node_type == "strong":
            # 太字
            children = node.get("children", [])
            for child in children:
                run = paragraph.add_run()
                run.text = self._node_to_text(child)
                run.font.bold = True
                run.font.size = Pt(18)
                self._apply_font_to_run(run)
                
        elif node_type == "emphasis":
            # 斜体
            children = node.get("children", [])
            for child in children:
                run = paragraph.add_run()
                run.text = self._node_to_text(child)
                run.font.italic = True
                run.font.size = Pt(18)
                self._apply_font_to_run(run)
                
        elif node_type == "codespan":
            # インラインコード - rawキーを優先的に使用
            run = paragraph.add_run()
            run.text = node.get("raw", "")
            run.font.name = "Consolas"  # コード用モノスペースフォント
            run.font.name_ascii = "Consolas"
            run.font.size = Pt(16)
            # 背景色を薄いグレーに設定
            run.font.fill.solid()
            run.font.fill.fore_color.rgb = RGBColor(60, 60, 60)  # 濃いめのグレー
            
        elif node_type in ["link", "image"]:
            # リンクは下線付きテキスト、画像は今後の課題として通常テキストで
            text = node.get("text", "") or node.get("raw", "")
            run = paragraph.add_run()
            run.text = text
            run.font.size = Pt(18)
            if node_type == "link":
                run.font.underline = True
            self._apply_font_to_run(run)
    
    def _extract_list_item_text(self, item: Dict[str, Any]) -> str:
        """リスト項目からテキストを効率的に抽出する改善版メソッド
        
        Args:
            item: リスト項目ノード
            
        Returns:
            str: 抽出されたテキスト
        """
        if not item or "children" not in item:
            return ""
            
        # リスト項目の子要素を取得
        children = item.get("children", [])
        result_text = ""
        
        # block_text要素を探す（最も一般的なリスト項目の構造）
        for child in children:
            if child.get("type") == "block_text" and "children" in child:
                block_children = child.get("children", [])
                
                # 子要素を順番に処理
                for block_child in block_children:
                    child_type = block_child.get("type", "")
                    
                    # テキスト要素
                    if child_type == "text":
                        result_text += block_child.get("raw", "")
                    
                    # コードスパン要素
                    elif child_type == "codespan":
                        result_text += block_child.get("raw", "")
                    
                    # 強調・太字要素
                    elif child_type in ["strong", "emphasis"]:
                        for grandchild in block_child.get("children", []):
                            result_text += grandchild.get("raw", "")
                            
                # block_text処理完了
                return result_text
        
        # その他のケース用のフォールバック処理
        for child in children:
            if "raw" in child:
                result_text += child.get("raw", "")
            elif "children" in child:
                # 再帰的に子要素を処理
                for grandchild in child.get("children", []):
                    if "raw" in grandchild:
                        result_text += grandchild.get("raw", "")
                    elif "children" in grandchild:
                        for gc in grandchild.get("children", []):
                            result_text += gc.get("raw", "")
        
        return result_text
    
    def build_presentation(self, slides_data: List[Dict[str, Any]], output_path: str) -> None:
        """スライドデータからプレゼンテーションを構築し保存する
        
        Args:
            slides_data: スライドデータのリスト
            output_path: 出力PPTXのパス
        """
        total_slides = len(slides_data)
        logger.info(f"{total_slides}枚のスライドを作成します")
        
        for slide_data in slides_data:
            self.create_slide(slide_data, total_slides)
        
        # 保存
        try:
            self.prs.save(output_path)
            logger.info(f"プレゼンテーションを保存しました: {output_path}")
        except Exception as e:
            logger.error(f"プレゼンテーション保存エラー: {e}")
            raise 